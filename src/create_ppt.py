# -*- coding: utf-8 -*-
import xlrd
from pptx import Presentation
from pptx.util import Inches, Pt


def highlight(workbook, prs):
    """ 读取网络工作表中的数据 """
    table = workbook.sheet_by_name(u'网络')
    for i in range(1, table.nrows):
        media, _, channel, title = table.row_values(i, 1, 5)
        if media == "":
            continue
        publish = xlrd.xldate_as_tuple(_, workbook.datemode)
        publish_date = "%d月%d日" % (publish[1], publish[2])
        highlight_ppt(media, channel, publish_date, title, prs)


def kol(workbook, prs):
    """ 读取KOL工作表中的数据 """
    table = workbook.sheet_by_name(u'KOL')
    for i in range(1, table.nrows):
        media, _, title = table.row_values(i, 3, 6)
        if media == "":
            continue
        publish = xlrd.xldate_as_tuple(_, workbook.datemode)
        publish_date = "%d月%d日" % (publish[1], publish[2])
        kol_ppt(media, publish_date, title, prs)


def app(workbook, prs):
    """ 读取APP工作表中的数据 """
    table = workbook.sheet_by_name(u'APP')
    for i in range(1, table.nrows):
        media, _, channel, title = table.row_values(i, 1, 5)
        if media == "":
            continue
        publish = xlrd.xldate_as_tuple(_, workbook.datemode)
        publish_date = "%d月%d日" % (publish[1], publish[2])
        app_ppt(media, publish_date, channel, title, prs)


def highlight_ppt(media, channel, publish_date, title, prs):
    """ 创建highlight模式PPT """
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    shape_title = slide.shapes.title
    shape_title.text = "Highlights"

    left = Inches(0.21)
    top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    tf.text = "媒体名称：" + media + "  频道：" + channel + "  发布时间：" + \
              publish_date + "\n" + "文章标题：" + title

    # left = Inches(0.66)
    # top = Inches(2)
    #
    # slide.shapes.add_picture("G:\\图片1.png", left, top)
    #
    # left = Inches(1.59)
    # top = Inches(6)
    # slide.shapes.add_picture("G:\\图片2.png", left, top)


def kol_ppt(media, publish_date, title, prs):
    """ 创建KOL模式PPT """
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    shape_title = slide.shapes.title
    shape_title.text = "KOL"

    # left = Inches(0.44)
    # top = Inches(1.31)
    # slide.shapes.add_picture("G:\\图片1.png", left, top)

    left = Inches(0.55)
    width = height = Inches(1)
    top = Inches(7.38)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    tf.text = "媒体名称：" + media + "  发布时间：" + \
              publish_date + "\n" + "文章标题：" + title


def app_ppt(media, publish_date, channel, title, prs):
    """ 创建APP模式PPT """
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    shape_title = slide.shapes.title
    shape_title.text = "APP"

    left = Inches(0.21)
    top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    tf.text = "媒体名称：" + media + "  频道：" + channel + "  发布时间：" + \
              publish_date + "\n" + "文章标题：" + title

    # left = Inches(2.49)
    # top = Inches(1.44)
    #
    # slide.shapes.add_picture("G:\\图片1.png", left, top)
    #
    # left = Inches(1.65)
    # top = Inches(5.58)
    # slide.shapes.add_picture("G:\\图片2.png", left, top)


if __name__ == "__main__":
    with open(".\\resource\\resource.txt", "r") as file:
        excel_file = file.readline().split("=")[1].strip("\n")
        template_ppt_file = file.readline().split("=")[1].strip("\n")
        output_ppt_file = file.readline().split("=")[1].strip("\n")
        workbook = xlrd.open_workbook(excel_file)
        prs = Presentation(template_ppt_file)
        kol(workbook, prs)
        app(workbook, prs)
        highlight(workbook, prs)
        prs.save(output_ppt_file)
